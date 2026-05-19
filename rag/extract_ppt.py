from pathlib import Path
from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parents[1]

PPT_PATH = BASE_DIR / "data" / "raw" / "OTML_Module1_completed.pptx"
OUTPUT_PATH = BASE_DIR / "data" / "processed" / "module1_otml_notes.txt"

def clean_text(text):
    """Clean unwanted encoding and spacing characters."""
    replacements = {
        "\xa0": " ",
        "Â": "",
        "\u200b": "",
        "\r": "\n",
    }

    for old, new in replacements.items():
        text = text.replace(old, new)

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)

def extract_text_from_shape(shape):
    """Extract text from normal text boxes and tables."""
    extracted_text = []

    if hasattr(shape, "text"):
        text = shape.text.strip()
        if text:
            extracted_text.append(text)

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_text.append(cell_text)
            if row_text:
                extracted_text.append(" | ".join(row_text))

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            extracted_text.extend(extract_text_from_shape(sub_shape))

    return extracted_text


def extract_ppt_text():
    if not PPT_PATH.exists():
        raise FileNotFoundError(f"PPT file not found: {PPT_PATH}")

    presentation = Presentation(PPT_PATH)

    all_slide_text = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        all_slide_text.append(f"\n\n===== Slide {slide_number} =====\n")

        slide_text = []

        for shape in slide.shapes:
            slide_text.extend(extract_text_from_shape(shape))

        if slide_text:
            all_slide_text.append("\n".join(slide_text))
        else:
            all_slide_text.append("[No readable text found on this slide]")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

    with open(OUTPUT_PATH, "w", encoding="utf-8") as file:
        file.write(clean_text("\n".join(all_slide_text)))

    print("PPT text extraction completed successfully.")
    print(f"Input file: {PPT_PATH}")
    print(f"Output file: {OUTPUT_PATH}")
    print(f"Total slides processed: {len(presentation.slides)}")


if __name__ == "__main__":
    extract_ppt_text()